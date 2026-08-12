from __future__ import annotations

import json
import shutil
from pathlib import Path
from typing import Iterable, TYPE_CHECKING

import pandas as pd
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

from .schemas import Document
from .security import safe_extract_zip, validate_upload, sanitize_filename
if TYPE_CHECKING:
    from .llm import GeminiGateway


class DocumentLoader:
    def __init__(self, ocr_gateway: GeminiGateway | None = None):
        self.ocr_gateway = ocr_gateway

    def expand_inputs(self, paths: Iterable[Path], workspace_dir: Path) -> list[Path]:
        expanded: list[Path] = []
        for path in paths:
            validate_upload(path)
            if path.suffix.lower() == ".zip":
                expanded.extend(safe_extract_zip(path, workspace_dir / "unzipped"))
            else:
                target = workspace_dir / "uploads" / sanitize_filename(path.name)
                target.parent.mkdir(parents=True, exist_ok=True)
                if path.resolve() != target.resolve():
                    shutil.copy2(path, target)
                expanded.append(target)
        return expanded

    def load(self, path: Path) -> tuple[list[Document], list[tuple[str, pd.DataFrame]]]:
        ext = path.suffix.lower()
        if ext == ".pdf":
            return self._pdf(path), []
        if ext == ".docx":
            return self._docx(path), []
        if ext == ".pptx":
            return self._pptx(path), []
        if ext == ".csv":
            df = pd.read_csv(path)
            return self._dataframe_docs(path.name, df), [(path.stem, df)]
        if ext in {".xlsx", ".xls"}:
            sheets = pd.read_excel(path, sheet_name=None)
            docs: list[Document] = []
            tables: list[tuple[str, pd.DataFrame]] = []
            for sheet, df in sheets.items():
                docs.extend(self._dataframe_docs(f"{path.name}:{sheet}", df))
                tables.append((f"{path.stem}_{sheet}", df))
            return docs, tables
        if ext == ".json":
            data = json.loads(path.read_text(encoding="utf-8", errors="ignore"))
            return [Document(json.dumps(data, indent=2, ensure_ascii=False), path.name)], []
        if ext in {".html", ".htm"}:
            soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "lxml")
            title = soup.title.string.strip() if soup.title and soup.title.string else None
            return [Document(soup.get_text("\n", strip=True), path.name, section=title)], []
        if ext in {".png", ".jpg", ".jpeg", ".webp"}:
            if not self.ocr_gateway:
                return [Document("[Image file indexed without OCR. Enable Gemini OCR to extract its text.]", path.name)], []
            return [Document(self.ocr_gateway.extract_file_text(path), path.name, metadata={"ocr": "gemini"})], []
        text = path.read_text(encoding="utf-8", errors="ignore")
        return [Document(text, path.name)], []

    def _pdf(self, path: Path) -> list[Document]:
        reader = PdfReader(str(path))
        docs: list[Document] = []
        total_chars = 0
        for i, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            total_chars += len(text)
            if text:
                docs.append(Document(text, path.name, page=i))
        if total_chars < 80 and self.ocr_gateway:
            extracted = self.ocr_gateway.extract_file_text(path)
            return [Document(extracted, path.name, metadata={"ocr": "gemini"})]
        return docs

    def _docx(self, path: Path) -> list[Document]:
        doc = DocxDocument(str(path))
        blocks: list[str] = []
        for p in doc.paragraphs:
            if p.text.strip():
                blocks.append(p.text.strip())
        for table in doc.tables:
            rows = []
            for row in table.rows:
                rows.append(" | ".join(cell.text.strip() for cell in row.cells))
            if rows:
                blocks.append("\n".join(rows))
        return [Document("\n\n".join(blocks), path.name)]

    def _pptx(self, path: Path) -> list[Document]:
        prs = Presentation(str(path))
        docs: list[Document] = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                docs.append(Document("\n".join(texts), path.name, page=i, section=f"Slide {i}"))
        return docs

    def _dataframe_docs(self, source: str, df: pd.DataFrame) -> list[Document]:
        docs: list[Document] = []
        clean = df.fillna("")
        for start in range(0, len(clean), 50):
            block = clean.iloc[start:start + 50]
            text = block.to_csv(index=False)
            docs.append(Document(text, source, section=f"Rows {start + 1}-{start + len(block)}", metadata={"structured": True}))
        if not docs:
            docs.append(Document("Columns: " + ", ".join(map(str, df.columns)), source, metadata={"structured": True}))
        return docs
